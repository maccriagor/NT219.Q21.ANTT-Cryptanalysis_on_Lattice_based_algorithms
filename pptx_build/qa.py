from pptx import Presentation
p = Presentation(r"../NT219_PQC_Report.pptx")
SW, SH = p.slide_width/914400, p.slide_height/914400
out = [f"slides={len(p.slides)}  size={SW:.2f}x{SH:.2f}in"]
for i, slide in enumerate(p.slides, 1):
    out.append(f"\n===== Slide {i} =====")
    for shp in slide.shapes:
        L = (shp.left or 0)/914400; T = (shp.top or 0)/914400
        Wd = (shp.width or 0)/914400; Hd = (shp.height or 0)/914400
        flags = []
        if L < -0.02 or T < -0.02: flags.append("OFF-NEG")
        if L+Wd > SW+0.05: flags.append(f"X-OVF({L+Wd:.2f})")
        if T+Hd > SH+0.05: flags.append(f"Y-OVF({T+Hd:.2f})")
        if shp.has_text_frame:
            txt = " / ".join(par.text for par in shp.text_frame.paragraphs if par.text.strip())
            if txt.strip():
                fl = (" <<"+",".join(flags)+">>") if flags else ""
                out.append(f"  ({L:.1f},{T:.1f} {Wd:.1f}x{Hd:.1f}) {txt[:130]}{fl}")
        elif flags:
            out.append(f"  (shape {L:.1f},{T:.1f} {Wd:.1f}x{Hd:.1f}) <<{','.join(flags)}>>")
open("qa_out.txt", "w", encoding="utf-8").write("\n".join(out))
print("OK -> qa_out.txt", len(out), "lines")
